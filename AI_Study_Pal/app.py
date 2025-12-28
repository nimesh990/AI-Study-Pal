import os
import sqlite3
from io import BytesIO
from datetime import datetime, timedelta
import csv
import random
import html
import requests

from flask import (
    Flask, render_template, request, session,
    send_file, redirect, url_for, g, flash
)

# NLTK
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize

# ICS
from icalendar import Calendar, Event

# PDF (ReportLab)
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.units import inch

# DOCX
from docx import Document

# PPTX
from pptx import Presentation
from pptx.util import Pt


app = Flask(__name__)
app.secret_key = "change-this-secret-key"

# ------------ Paths ------------
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
DATA_FILE = os.path.join(BASE_DIR, "study_logs.csv")
DB_PATH = os.path.join(BASE_DIR, "study_pal.db")

CSV_HEADER = ["subject", "hours", "start_date", "start_time", "days", "break_min", "weekly_off"]


# ======================================================
#                 SQLITE HELPERS (USERS + PLANS)
# ======================================================

def get_db():
    db = getattr(g, "_db", None)
    if db is None:
        db = g._db = sqlite3.connect(DB_PATH)
        db.row_factory = sqlite3.Row
    return db


@app.teardown_appcontext
def close_db(exception):
    db = getattr(g, "_db", None)
    if db is not None:
        db.close()


def init_db():
    conn = sqlite3.connect(DB_PATH)
    cur = conn.cursor()

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL
        );
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS plans (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL,
            subject TEXT NOT NULL,
            hours INTEGER NOT NULL,
            start_date TEXT NOT NULL,
            start_time TEXT NOT NULL,
            days INTEGER NOT NULL,
            break_min INTEGER NOT NULL,
            weekly_off TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY(user_id) REFERENCES users(id)
        );
        """
    )

    conn.commit()
    conn.close()


init_db()


def get_current_user_id():
    return session.get("user_id")


def get_current_username():
    return session.get("username")


# ======================================================
#                 CSV COMPATIBILITY (OLD DATA)
# ======================================================

def ensure_csv_schema():
    if not os.path.exists(DATA_FILE):
        with open(DATA_FILE, mode="w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(CSV_HEADER)
        return

    with open(DATA_FILE, mode="r", newline="", encoding="utf-8-sig") as f:
        raw_lines = [line.strip() for line in f if line.strip()]

    if not raw_lines:
        with open(DATA_FILE, mode="w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(CSV_HEADER)
        return

    first = [x.strip() for x in raw_lines[0].split(",")]
    already_header = (first == CSV_HEADER)

    data_lines = raw_lines[1:] if already_header else raw_lines
    normalized_rows = []

    for line in data_lines:
        parts = [p.strip() for p in line.split(",")]
        if not parts or all(p == "" for p in parts):
            continue

        if len(parts) == 2:
            subject, hours = parts
            normalized_rows.append([subject, hours, "", "", "", "0", "none"])
        elif len(parts) == 5:
            subject, hours, start_date, start_time, days = parts
            normalized_rows.append([subject, hours, start_date, start_time, days, "0", "none"])
        else:
            row7 = (parts + ["", "", "", "0", "none"])[:7]
            if not row7[6]:
                row7[6] = "none"
            if not row7[5]:
                row7[5] = "0"
            normalized_rows.append(row7)

    bak = DATA_FILE + ".bak"
    try:
        if os.path.exists(bak):
            os.remove(bak)
        os.rename(DATA_FILE, bak)
    except Exception:
        pass

    with open(DATA_FILE, mode="w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(CSV_HEADER)
        writer.writerows(normalized_rows)


ensure_csv_schema()


# ======================================================
#                 PLAN HELPERS
# ======================================================

def weekday_code_to_index(code: str):
    mapping = {"mon": 0, "tue": 1, "wed": 2, "thu": 3, "fri": 4, "sat": 5, "sun": 6}
    return mapping.get(code)


def build_daily_blocks(hours: int, break_min: int, start_dt: datetime):
    total_min = max(1, hours * 60)
    break_min = max(0, break_min)

    if break_min == 0 or break_min >= total_min:
        return [("Study Session", start_dt, start_dt + timedelta(minutes=total_min))]

    first_min = total_min // 2
    second_min = total_min - first_min

    s1_start = start_dt
    s1_end = s1_start + timedelta(minutes=first_min)

    s2_start = s1_end + timedelta(minutes=break_min)
    s2_end = s2_start + timedelta(minutes=second_min)

    return [("Study Session 1", s1_start, s1_end), ("Study Session 2", s2_start, s2_end)]


def get_last_plan_or_400():
    if "last_schedule" not in session:
        return None, ("No schedule found. Generate a study plan first.", 400)

    subject = session.get("last_subject", "Study")
    hours = int(session.get("last_hours", 1))
    days = int(session.get("last_days", 7))
    start_date_str = session.get("last_start_date")
    start_time_str = session.get("last_start_time")
    break_min = int(session.get("last_break_min", 0))
    weekly_off = session.get("last_weekly_off", "none")
    schedule = session.get("last_schedule", [])
    return (subject, hours, days, start_date_str, start_time_str, break_min, weekly_off, schedule), None


# ======================================================
#                     AUTH ROUTES
# ======================================================

@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "GET":
        return render_template("register.html")

    username = request.form.get("username", "").strip()
    password = request.form.get("password", "").strip()

    if not username or not password:
        flash("Username and password required.")
        return render_template("register.html")

    db = get_db()
    try:
        db.execute("INSERT INTO users (username, password) VALUES (?, ?)", (username, password))
        db.commit()
    except sqlite3.IntegrityError:
        flash("Username already taken.")
        return render_template("register.html")

    flash("Registered successfully. Please log in.")
    return redirect(url_for("login"))


@app.route("/login", methods=["GET", "POST"])
def login():
    if request.method == "GET":
        return render_template("login.html")

    username = request.form.get("username", "").strip()
    password = request.form.get("password", "").strip()

    db = get_db()
    row = db.execute(
        "SELECT id, username, password FROM users WHERE username = ?",
        (username,)
    ).fetchone()

    if row is None or row["password"] != password:
        flash("Invalid username or password.")
        return render_template("login.html")

    session["user_id"] = row["id"]
    session["username"] = row["username"]
    flash("Logged in successfully.")
    return redirect(url_for("home"))


@app.route("/logout")
def logout():
    session.pop("user_id", None)
    session.pop("username", None)
    flash("Logged out.")
    return redirect(url_for("home"))


@app.route("/my_plans")
def my_plans():
    user_id = get_current_user_id()
    if not user_id:
        flash("Please log in to see your plans.")
        return redirect(url_for("login"))

    db = get_db()
    rows = db.execute(
        """
        SELECT id, subject, hours, start_date, start_time, days,
               break_min, weekly_off, created_at
        FROM plans
        WHERE user_id = ?
        ORDER BY created_at DESC
        """,
        (user_id,)
    ).fetchall()

    return render_template("my_plans.html", plans=rows)


# ======================================================
#                     MAIN PAGES
# ======================================================

@app.route("/")
def home():
    today = datetime.now().date().strftime("%Y-%m-%d")
    return render_template("home.html", default_date=today, username=get_current_username())


@app.route("/generate", methods=["POST"])
def generate():
    subject = request.form.get("subject", "").strip()
    hours = int(request.form.get("hours", 1))
    start_date_str = request.form.get("start_date", "").strip()
    start_time_str = request.form.get("start_time", "").strip()
    days = int(request.form.get("days", 7))
    break_min = int(request.form.get("break_min", 0))
    weekly_off = request.form.get("weekly_off", "none").strip().lower()

    if not subject:
        return "Subject is required.", 400
    if days < 1 or days > 60:
        return "Days must be between 1 and 60.", 400
    if hours < 1 or hours > 12:
        return "Hours per day must be between 1 and 12.", 400
    if break_min < 0 or break_min > 180:
        return "Break minutes must be between 0 and 180.", 400

    try:
        start_date = datetime.strptime(start_date_str, "%Y-%m-%d").date()
    except Exception:
        return "Invalid start date. Use YYYY-MM-DD.", 400

    try:
        hh, mm = start_time_str.split(":")
        start_hour = int(hh)
        start_minute = int(mm)
        if not (0 <= start_hour <= 23 and 0 <= start_minute <= 59):
            return "Invalid start time.", 400
    except Exception:
        return "Invalid start time. Use HH:MM.", 400

    off_index = None
    if weekly_off != "none":
        off_index = weekday_code_to_index(weekly_off)
        if off_index is None:
            return "Invalid weekly off value.", 400

    # CSV log (for stats)
    with open(DATA_FILE, mode="a", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow([subject, hours, start_date_str, start_time_str, days, break_min, weekly_off])

    # Build schedule
    schedule = []
    for i in range(days):
        d = start_date + timedelta(days=i)

        if off_index is not None and d.weekday() == off_index:
            schedule.append(f"Day {i+1} ({d.strftime('%Y-%m-%d')}): WEEKLY OFF")
            continue

        start_dt = datetime(d.year, d.month, d.day, start_hour, start_minute, 0)
        blocks = build_daily_blocks(hours, break_min, start_dt)

        if len(blocks) == 1:
            _, b1s, b1e = blocks[0]
            schedule.append(
                f"Day {i+1} ({d.strftime('%Y-%m-%d')}): Study {subject} "
                f"{b1s.strftime('%H:%M')} - {b1e.strftime('%H:%M')} ({hours}h)"
            )
        else:
            _, b1s, b1e = blocks[0]
            _, b2s, b2e = blocks[1]
            schedule.append(
                f"Day {i+1} ({d.strftime('%Y-%m-%d')}): "
                f"{subject} {b1s.strftime('%H:%M')}-{b1e.strftime('%H:%M')}, "
                f"Break {break_min} min, "
                f"{subject} {b2s.strftime('%H:%M')}-{b2e.strftime('%H:%M')}"
            )

    weekly_off_text = weekly_off.upper() if weekly_off != "none" else "NONE"
    plan = (
        f"Study {subject} for {hours} hours/day for {days} days starting {start_date_str} at {start_time_str}. "
        f"Break: {break_min} min. Weekly off: {weekly_off_text}."
    )

    # Save plan in session for downloads
    session["last_subject"] = subject
    session["last_hours"] = hours
    session["last_days"] = days
    session["last_start_date"] = start_date_str
    session["last_start_time"] = start_time_str
    session["last_break_min"] = break_min
    session["last_weekly_off"] = weekly_off
    session["last_schedule"] = schedule

    # Save plan in DB if logged in
    user_id = get_current_user_id()
    if user_id:
        db = get_db()
        db.execute(
            """
            INSERT INTO plans (
                user_id, subject, hours, start_date, start_time,
                days, break_min, weekly_off, created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                user_id,
                subject,
                hours,
                start_date_str,
                start_time_str,
                days,
                break_min,
                weekly_off,
                datetime.now().isoformat(timespec="seconds"),
            ),
        )
        db.commit()

    return render_template("result.html", plan=plan, schedule=schedule, username=get_current_username())


# ---------------- DOWNLOAD ROUTES (ICS/PDF/DOCX/PPTX) ----------------
@app.route("/download_ics")
def download_ics():
    data, err = get_last_plan_or_400()
    if err:
        return err

    subject, hours, days, start_date_str, start_time_str, break_min, weekly_off, schedule = data
    start_date = datetime.strptime(start_date_str, "%Y-%m-%d").date()
    hh, mm = start_time_str.split(":")
    start_hour = int(hh)
    start_minute = int(mm)

    off_index = None
    if weekly_off != "none":
        off_index = weekday_code_to_index(weekly_off)

    cal = Calendar()
    cal.add("prodid", "-//AI Study Pal//")
    cal.add("version", "2.0")

    for i in range(days):
        day_date = start_date + timedelta(days=i)
        if off_index is not None and day_date.weekday() == off_index:
            continue

        start_dt = datetime(day_date.year, day_date.month, day_date.day, start_hour, start_minute, 0)
        blocks = build_daily_blocks(hours, break_min, start_dt)

        for title, dt_start, dt_end in blocks:
            event = Event()
            event.add("summary", f"Study {subject} (Day {i+1}) - {title}")
            event.add("dtstart", dt_start)
            event.add("dtend", dt_end)
            event.add("description", f"AI Study Pal: {subject} ({title}).")
            cal.add_component(event)

    bio = BytesIO(cal.to_ical())
    bio.seek(0)
    filename = f"study_plan_{subject.replace(' ', '_')}.ics"
    return send_file(bio, as_attachment=True, download_name=filename, mimetype="text/calendar")


@app.route("/download_pdf")
def download_pdf():
    data, err = get_last_plan_or_400()
    if err:
        return err

    subject, hours, days, start_date_str, start_time_str, break_min, weekly_off, schedule = data

    pdf_buffer = BytesIO()
    doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)

    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title", parent=styles["Heading1"], fontSize=20, spaceAfter=16, alignment=1)
    heading_style = ParagraphStyle("Heading", parent=styles["Heading2"], fontSize=13, spaceAfter=10, spaceBefore=10)

    content = []
    content.append(Paragraph("AI Study Pal - Study Schedule", title_style))
    content.append(Spacer(1, 0.2 * inch))
    content.append(Paragraph(f"Subject: <b>{subject}</b>", styles["Normal"]))
    content.append(Paragraph(f"Hours per day: <b>{hours}</b>", styles["Normal"]))
    content.append(Paragraph(f"Days: <b>{days}</b>", styles["Normal"]))
    content.append(Paragraph(f"Start: <b>{start_date_str} {start_time_str}</b>", styles["Normal"]))
    content.append(Paragraph(f"Break: <b>{break_min} min</b>", styles["Normal"]))
    content.append(Paragraph(f"Weekly off: <b>{weekly_off.upper()}</b>", styles["Normal"]))
    content.append(Spacer(1, 0.2 * inch))

    content.append(Paragraph("Schedule", heading_style))
    for line in schedule:
        content.append(Paragraph(f"• {line}", styles["Normal"]))

    doc.build(content)
    pdf_buffer.seek(0)

    filename = f"study_plan_{subject.replace(' ', '_')}.pdf"
    return send_file(pdf_buffer, as_attachment=True, download_name=filename, mimetype="application/pdf")


@app.route("/download_docx")
def download_docx():
    data, err = get_last_plan_or_400()
    if err:
        return err

    subject, hours, days, start_date_str, start_time_str, break_min, weekly_off, schedule = data

    doc = Document()
    doc.add_heading("AI Study Pal - Study Schedule", level=1)
    doc.add_paragraph(f"Subject: {subject}")
    doc.add_paragraph(f"Hours per day: {hours}")
    doc.add_paragraph(f"Days: {days}")
    doc.add_paragraph(f"Start: {start_date_str} {start_time_str}")
    doc.add_paragraph(f"Break: {break_min} min")
    doc.add_paragraph(f"Weekly off: {weekly_off.upper()}")

    doc.add_heading("Schedule", level=2)
    for line in schedule:
        doc.add_paragraph(line, style="List Bullet")

    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)

    filename = f"study_plan_{subject.replace(' ', '_')}.docx"
    return send_file(
        bio,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@app.route("/download_pptx")
def download_pptx():
    data, err = get_last_plan_or_400()
    if err:
        return err

    subject, hours, days, start_date_str, start_time_str, break_min, weekly_off, schedule = data

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Study Pal"
    slide.placeholders[1].text = f"{subject} | {hours}h/day | {days} days | Break {break_min}m | Off {weekly_off.upper()}"

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"Schedule (Start: {start_date_str} {start_time_str})"
    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, line in enumerate(schedule[:20]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(16)

    if len(schedule) > 20:
        p = tf.add_paragraph()
        p.text = f"...and {len(schedule) - 20} more lines"
        p.level = 0
        p.font.size = Pt(16)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    filename = f"study_plan_{subject.replace(' ', '_')}.pptx"
    return send_file(
        bio,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )


# ---------------- STATS / SUMMARY / TIPS / RESOURCES / QUIZ (same as before) ----------------
@app.route("/stats")
def stats():
    subjects = {}
    total_plans = 0

    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, mode="r", newline="", encoding="utf-8-sig") as f:
            lines = (line for line in f if line.strip())
            reader = csv.DictReader(lines)
            for row in reader:
                subj_raw = row.get("subject")
                if not subj_raw:
                    continue
                total_plans += 1
                subj = subj_raw.strip().lower()
                subjects[subj] = subjects.get(subj, 0) + 1

    subject_counts = sorted(subjects.items(), key=lambda x: x[1], reverse=True)
    return render_template("stats.html", total_plans=total_plans, subject_counts=subject_counts)


@app.route("/summary", methods=["GET", "POST"])
def summary():
    if request.method == "GET":
        return render_template("summary.html", original_text="", summarized_text="")

    text = request.form["text"]
    sentences = [s.strip() for s in text.split(".") if s.strip()]
    short = ". ".join(sentences[:2]) if len(sentences) > 2 else ". ".join(sentences)
    if short and not short.endswith("."):
        short += "."
    return render_template("summary.html", original_text=text, summarized_text=short)


@app.route("/tips", methods=["GET", "POST"])
def tips():
    if request.method == "GET":
        return render_template("tips.html", text="", tips=[])

    text = request.form["text"]
    cleaned = text.lower()
    tokens = word_tokenize(cleaned)

    default_stops = set(stopwords.words("english"))
    custom_stops = {"today", "yesterday", "tomorrow", "studied", "study", "studying", "chapter", "topic", "lesson"}
    all_stops = default_stops.union(custom_stops)

    words = [w for w in tokens if w.isalpha() and w not in all_stops and len(w) > 3]
    seen = set()
    uniq = []
    for w in words:
        if w not in seen:
            seen.add(w)
            uniq.append(w)

    top = uniq[:5]
    tips_list = []
    for kw in top:
        tips_list.append(f"Revise the concept of '{kw}' regularly.")
        tips_list.append(f"Practice 2-3 questions daily related to '{kw}'.")

    if not tips_list:
        tips_list.append("Text is too short. Add more content to get useful tips.")

    return render_template("tips.html", text=text, tips=tips_list)


@app.route("/resources", methods=["GET", "POST"])
def resources():
    suggestions = []
    if request.method == "POST":
        subject = request.form["subject"].strip().lower()

        resource_map = {
            "math": ["https://www.khanacademy.org/math", "https://www.cuemath.com/algebra/", "https://www.mathsisfun.com/"],
            "physics": ["https://www.khanacademy.org/science/physics", "https://www.physicsclassroom.com/", "https://www.hyperphysics.phy-astr.gsu.edu/"],
            "chemistry": ["https://www.khanacademy.org/science/chemistry", "https://chem.libretexts.org/", "https://www.chemguide.co.uk/"],
            "computer science": ["https://www.geeksforgeeks.org/", "https://www.w3schools.com/", "https://www.programiz.com/python-programming"]
        }

        if subject in ["maths", "mathematics"]:
            subject_key = "math"
        elif subject in ["cs", "computer", "computer science", "cse"]:
            subject_key = "computer science"
        else:
            subject_key = subject

        suggestions = resource_map.get(subject_key, [])
        if not suggestions:
            suggestions = ["No specific resources found for this subject.", "Try searching on Google or YouTube with detailed topic names."]

    return render_template("resources.html", suggestions=suggestions)


SUBJECT_MAP = {
    "gk": 9,
    "general knowledge": 9,
    "computer": 18,
    "cs": 18,
    "science": 17,
    "history": 23
}


def fetch_questions_from_api(subject, amount=5):
    category_id = SUBJECT_MAP.get(subject.lower(), 9)
    url = "https://opentdb.com/api.php"
    params = {"amount": amount, "category": category_id, "type": "multiple"}

    r = requests.get(url, params=params)
    data = r.json()

    questions = []
    if data.get("response_code") != 0:
        return questions

    for i, item in enumerate(data["results"], start=1):
        q_text = html.unescape(item["question"])
        correct = html.unescape(item["correct_answer"])
        incorrect = [html.unescape(x) for x in item["incorrect_answers"]]
        options = incorrect + [correct]
        random.shuffle(options)
        questions.append({"id": f"q{i}", "q": q_text, "options": options, "answer": correct})

    return questions


@app.route("/api_quiz", methods=["GET", "POST"])
def api_quiz():
    if request.method == "GET":
        return render_template("api_select_subject.html")

    subject = request.form["subject"]
    questions = fetch_questions_from_api(subject, amount=5)
    if not questions:
        return render_template("api_select_subject.html", error="Unable to fetch questions from API. Please try again.")
    return render_template("api_quiz.html", subject=subject, questions=questions)


@app.route("/api_quiz_result", methods=["POST"])
def api_quiz_result():
    subject = request.form["subject"]
    total = int(request.form["total"])

    questions = []
    score = 0
    for i in range(1, total + 1):
        qid = f"q{i}"
        q_text = request.form.get(f"text_{qid}")
        correct = request.form.get(f"ans_{qid}")
        given = request.form.get(qid)
        if given == correct:
            score += 1
        questions.append({"question": q_text, "given": given, "correct": correct})

    return render_template("api_result.html", subject=subject, score=score, total=total, questions=questions)


if __name__ == "__main__":
    app.run(debug=True)
